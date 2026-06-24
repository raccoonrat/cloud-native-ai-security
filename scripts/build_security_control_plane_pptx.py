#!/usr/bin/env python3
"""Generate editable Chinese .pptx for Cloud Native Enterprise AI Security Control Plane v1.2."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "presentations" / "Cloud-Native-Enterprise-AI-Security-Control-Plane-v1.2.pptx"

# Teal Trust palette (security / enterprise)
C_PRIMARY = RGBColor(0x06, 0x5A, 0x82)
C_ACCENT = RGBColor(0x02, 0xC3, 0x9A)
C_DARK = RGBColor(0x21, 0x29, 0x5C)
C_INK = RGBColor(0x1E, 0x27, 0x61)
C_MUTED = RGBColor(0x5A, 0x6B, 0x7D)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)

FONT = "Microsoft YaHei"


def set_slide_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_footer_bar(slide, text: str):
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(5.15), Inches(10), Inches(0.475)
    )  # MSO_SHAPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_PRIMARY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = C_WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT


def style_title(tf, size=32, color=C_INK, bold=True):
    tf.word_wrap = True
    for i, p in enumerate(tf.paragraphs):
        p.font.name = FONT
        p.font.size = Pt(size if i == 0 else size - 6)
        p.font.bold = bold and i == 0
        p.font.color.rgb = color
        p.space_after = Pt(8)


def add_bullets(tf, items, size=18, color=C_INK, level0=True):
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.level = 0
        p.space_after = Pt(10)
        if level0:
            p.bullet = True


def add_title_content_slide(prs, title, bullets, footer, dark=False):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    if dark:
        set_slide_bg(slide, C_DARK)
        title_color, body_color = C_WHITE, RGBColor(0xE8, 0xED, 0xF5)
    else:
        set_slide_bg(slide, C_LIGHT)
        title_color, body_color = C_INK, C_INK

    slide.shapes.title.text = title
    style_title(slide.shapes.title.text_frame, size=30, color=title_color)

    body = slide.placeholders[1].text_frame
    add_bullets(body, bullets, size=17, color=body_color)
    add_footer_bar(slide, footer)
    return slide


def add_section_slide(prs, section, subtitle, footer):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    set_slide_bg(slide, C_PRIMARY)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    tf.text = section
    p = tf.paragraphs[0]
    p.font.name = FONT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C_WHITE
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.name = FONT
    p2.font.size = Pt(20)
    p2.font.color.rgb = C_ACCENT
    p2.space_before = Pt(16)
    add_footer_bar(slide, footer)
    return slide


def add_stat_slide(prs, title, stats, footer):
    """stats: list of (label, value, note)"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, C_LIGHT)
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.45), Inches(9), Inches(0.8))
    tb.text_frame.text = title
    style_title(tb.text_frame, size=28, color=C_INK)

    cols = min(len(stats), 3)
    w = 9.0 / cols
    for i, (label, value, note) in enumerate(stats):
        col = i % cols
        row = i // cols
        x = 0.6 + col * w
        y = 1.5 + row * 1.85
        card = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w - 0.25), Inches(1.55))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = RGBColor(0xD0, 0xD7, 0xDE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.text = label
        p0.font.name = FONT
        p0.font.size = Pt(11)
        p0.font.color.rgb = C_MUTED
        p1 = tf.add_paragraph()
        p1.text = value
        p1.font.name = FONT
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = C_PRIMARY
        p1.space_before = Pt(4)
        p2 = tf.add_paragraph()
        p2.text = note
        p2.font.name = FONT
        p2.font.size = Pt(10)
        p2.font.color.rgb = C_MUTED
        p2.space_before = Pt(2)
    add_footer_bar(slide, footer)
    return slide


def build():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # S01 Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(slide, C_DARK)
    slide.shapes.title.text = "云原生企业 AI 安全控制平面"
    sub = slide.placeholders[1]
    sub.text = "Cloud Native Enterprise AI Security Control Plane\nV1.2 · MCP / Tool-Calling Runtime Security Hardening"
    for shape in (slide.shapes.title, sub):
        tf = shape.text_frame
        for p in tf.paragraphs:
            p.font.name = FONT
            p.font.color.rgb = C_WHITE
            p.font.size = Pt(36 if shape == slide.shapes.title else 16)
    add_footer_bar(slide, "工程 / 平台 / 安全 / SRE · 2026")

    # S02 Executive Summary
    add_title_content_slide(
        prs,
        "执行摘要",
        [
            "统一 AI 运行时安全控制平面：input、RAG、tool、agent、output 全链路",
            "策略驱动：policy runtime 按租户/区域/角色/数据等级执行安全行为",
            "多信号融合：detector 产信号，decision engine 合成最终动作",
            "可解释决策：每次 allow/block/redact 含 reason_code + matched rules + evidence",
            "审计级追溯：decision trace 支持合规、复盘、offline replay、release gate",
            "云原生可扩展：K8s / sidecar / gateway / SDK，多租户多区域灰度",
        ],
        "01 / 18 · 执行摘要",
    )

    # S03 Shift
    add_title_content_slide(
        prs,
        "范式转移：从文本生成到工具执行",
        [
            "Enterprise AI 进入运维控制面：MCP / tool-calling 带来真实副作用",
            "Northbound API 认证、TLS、限流仍然必要，但已不充分",
            "自然语言意图可能被错误翻译为 API 调用；prompt injection 可影响工具与参数",
            "这些问题应抽象为平台级能力，而非某产品 MCP 集成的局部补丁",
        ],
        "02 / 18 · 范式转移",
    )

    # S04 Risks
    add_stat_slide(
        prs,
        "系统解决的核心风险",
        [
            ("Injection", "Prompt + RAG", "注入与检索污染"),
            ("Leakage", "PII / 机密", "输入输出泄密"),
            ("Tool", "Abuse", "越权与危险参数"),
            ("Automation", "Agent", "多步自主失控"),
            ("Transport", "Trust", "TLS / 证书 / 确认"),
            ("Evidence", "Gap", "无法证明决策依据"),
        ],
        "03 / 18 · 风险全景",
    )

    # S05 Scope
    add_title_content_slide(
        prs,
        "系统定位：是什么 / 不是什么",
        [
            "✓ 企业 AI 运行时安全控制平面（非单点 detector / prompt filter）",
            "✓ 策略驱动 + 上下文感知 + 多信号融合 + 可解释决策",
            "✗ 非模型训练安全平台、非 IAM 替代品、非完整 DLP 替代",
            "✗ 非单一内容审核服务、非仅依赖 LLM Judge 的 guardrail",
        ],
        "04 / 18 · 系统定位",
    )

    add_section_slide(prs, "第一幕", "架构原则与控制平面", "05 / 18")

    # S07 Control Plane
    add_title_content_slide(
        prs,
        "设计原则：Control Plane First",
        [
            "统一流水线：Context → Signals → Policy → Decision → Enforcement",
            "→ Trace → Evaluation → Release Gate → Feedback",
            "实现统一策略治理、审计证据、发布门禁、风险度量、回放复现",
            "拒绝散落在业务代码中的 if/else 安全判断",
        ],
        "06 / 18 · 控制平面优先",
        dark=True,
    )

    # S08 Detector
    add_title_content_slide(
        prs,
        "设计原则：Detector Does Not Decide",
        [
            "Detector 只产生 normalized signal，不直接 block/allow",
            "Decision Engine 负责：signal fusion + policy matching + action synthesis",
            "每次决策携带 reason_code、matched_rules、evidence 引用",
            "Detector 版本变更纳入 replay key，保证决策可复现",
        ],
        "07 / 18 · 检测与决策分离",
    )

    # S09 Architecture
    add_title_content_slide(
        prs,
        "八层架构",
        [
            "① Enterprise Request  ② Context Builder  ③ Detection Layer",
            "④ Signal Normalization & Fusion  ⑤ Policy Runtime",
            "⑥ Decision Engine  ⑦ Enforcement Layer",
            "⑧ Trace / Evidence / Evaluation / Release Gate",
        ],
        "08 / 18 · 高层架构",
    )

    add_section_slide(prs, "第二幕", "MCP / Tool-Calling 安全边界", "09 / 18")

    # S11 MCP
    add_title_content_slide(
        prs,
        "MCP 是 AI 中介执行边界",
        [
            "MCP / tool-calling ≠ 对既有 API 的 thin wrapper",
            "所有 tool-mediated 操作必须在 tool_pre_execution / agent_step_check 过控制平面",
            "绑定 user → client → MCP server → tool → API → resource → action",
            "通用 API wrapper 工具默认 block，除非 policy 显式批准",
            "disruptive 操作须 explicit_user_confirmation + confirmation_challenge",
        ],
        "10 / 18 · MCP 边界",
        dark=True,
    )

    # S12 Trust + Separation
    add_title_content_slide(
        prs,
        "信任链与指令/数据分离",
        [
            "Tool-to-API Trust Chain：每一跳绑定身份，参数 schema + resource-scope 校验",
            "Instruction/Data Separation：用户 prompt、RAG chunk、tool output = 不可信数据",
            "系统 prompt、policy snapshot、批准 schema、用户确认 = 可信指令",
            "外部内容不得成为 agent loop 中的可执行指令",
        ],
        "11 / 18 · 信任链",
    )

    add_section_slide(prs, "第三幕", "运行时契约与 Enforcement", "12 / 18")

    # S14 Contracts
    add_title_content_slide(
        prs,
        "Canonical Contract（硬约束源）",
        [
            "stage 枚举：input_precheck · rag_* · tool_pre/post · agent_* · model_output_check",
            "request_type 已 deprecated，禁止用于策略/trace/eval 匹配",
            "Action 状态机：terminal / transform / human-async / observability",
            "require_approval 是 pending state，须携带 approval_ticket_id",
            "Replay：content hash + context hash + detector/policy/engine versions",
        ],
        "13 / 18 · 运行时契约",
    )

    # S15 Enforcement
    add_stat_slide(
        prs,
        "四类 Enforcement 控制面",
        [
            ("Input", "input_precheck", "注入 · PII · 策略绑定"),
            ("RAG", "chunk check", "来源信任 · remove_chunk"),
            ("Tool", "pre_execution", "V1.2 重点 · 参数校验"),
            ("Output", "output_check", "泄密 · redact · block"),
        ],
        "14 / 18 · Enforcement",
    )

    # S16 Policy Trace
    add_title_content_slide(
        prs,
        "Policy DSL · Trace · Release Gate",
        [
            "Policy DSL v0.1：validate / lint / diff，immutable policy snapshot",
            "两阶段 Trace：同步 minimal commit（trace_id 100%）+ 异步 evidence 补充",
            "Release Gate：detector/policy/schema 发布阻断指标回归",
            "enforcement_mode：monitor · shadow · active · canary",
        ],
        "15 / 18 · 治理与证据",
        dark=True,
    )

    # S17 Roadmap
    add_stat_slide(
        prs,
        "工程实施路线",
        [
            ("P0", "Vertical Slice", "input + output 闭环 · Replay CLI"),
            ("P1", "Policy Hardening", "DSL · Evidence · Gate basic"),
            ("P1.5", "MCP Slice", "tool_pre_execution 最小闭环"),
            ("P2", "RAG + Agent", "chunk 安全 · agent step"),
            ("P3", "Production", "online sampling · auto rollback"),
        ],
        "16 / 18 · 路线图",
    )

    # S18 Close
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_slide_bg(slide, C_DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(2.8))
    tf = box.text_frame
    tf.text = "核心判断"
    p = tf.paragraphs[0]
    p.font.name = FONT
    p.font.size = Pt(22)
    p.font.color.rgb = C_ACCENT
    q = tf.add_paragraph()
    q.text = "Northbound API 安全是 trust chain 的必要一环，\nMCP/tool 层必须独立做意图、参数、确认与审计控制。"
    q.font.name = FONT
    q.font.size = Pt(26)
    q.font.bold = True
    q.font.color.rgb = C_WHITE
    q.space_before = Pt(20)
    q2 = tf.add_paragraph()
    q2.text = "先跑通 P0 vertical slice，再 ship P1.5 给 MCP 集成场景。"
    q2.font.name = FONT
    q2.font.size = Pt(16)
    q2.font.color.rgb = RGBColor(0xC0, 0xC8, 0xD4)
    q2.space_before = Pt(24)
    add_footer_bar(slide, "17 / 18 · 谢谢")

    prs.save(str(OUT))
    print(f"Saved: {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
